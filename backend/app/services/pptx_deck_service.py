"""
PPTX Deck Service — Data Quantum theme preset
============================================
Render engine desacoplado del LLM:
 - El LLM aporta contenido semantico (titulo, bullets, script).
 - Este servicio aplica SIEMPRE un layout corporativo fijo.

Layouts soportados:
 - cover
 - agenda
 - section-divider
 - content-bullets
 - content-two-columns
 - closing / next-steps
"""

from __future__ import annotations

import json
import os
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple


def _static_generated_root() -> str:
    app_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # backend/app
    static_dir = os.path.join(app_dir, "static", "generated-courses")
    os.makedirs(static_dir, exist_ok=True)
    return static_dir


def _theme_config_path() -> str:
    return os.path.join(os.path.dirname(os.path.abspath(__file__)), "pptx_theme_data_quantum.json")


def _rgb_from_hex(hex_value: str):
    from pptx.dml.color import RGBColor

    h = (hex_value or "").strip().lstrip("#")
    if len(h) != 6:
        h = "000000"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


@dataclass
class DeckTheme:
    theme_name: str
    palette: Dict[str, str]
    typography: Dict[str, float]
    spacing: Dict[str, float]
    limits: Dict[str, int]
    structure: Dict[str, int]
    layout_order: List[str]


def _load_theme() -> DeckTheme:
    path = _theme_config_path()
    with open(path, "r", encoding="utf-8") as f:
        cfg = json.load(f)
    return DeckTheme(
        theme_name=cfg.get("theme_name", "Data Quantum Corporate"),
        palette=cfg.get("palette", {}),
        typography=cfg.get("typography", {}),
        spacing=cfg.get("spacing", {}),
        limits=cfg.get("limits", {}),
        structure=cfg.get("structure", {}),
        layout_order=cfg.get("layout_order", []),
    )


def _safe_title(slide: Dict[str, Any], idx: int) -> str:
    title = (slide.get("title") or "").strip()
    return title if title else f"Slide {idx}"


def _split_bullets(
    bullets: List[str],
    *,
    max_bullets: int,
    max_chars: int,
) -> Tuple[List[List[str]], int]:
    """
    Normaliza bullets en chunks sin sobrecargar la slide.
    Devuelve (chunks, num_truncados).
    """
    if not bullets:
        return [[]], 0

    normalized: List[str] = []
    trunc_count = 0
    for b in bullets:
        text = " ".join(str(b or "").split())
        if len(text) > max_chars:
            text = text[: max_chars - 1].rstrip() + "…"
            trunc_count += 1
        if text:
            normalized.append(text)

    if not normalized:
        return [[]], trunc_count

    chunks: List[List[str]] = []
    for i in range(0, len(normalized), max_bullets):
        chunks.append(normalized[i : i + max_bullets])
    return chunks, trunc_count


def _add_background(slide_obj, *, width, height, color_hex: str):
    from pptx.util import Inches

    bg = slide_obj.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(0),
        Inches(0),
        width,
        height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb_from_hex(color_hex)
    bg.line.fill.background()


def _apply_font(run_or_para, *, font_name: str, size_pt: int, color_hex: str, bold: bool = False):
    from pptx.util import Pt

    if hasattr(run_or_para, "font"):
        f = run_or_para.font
    else:
        f = run_or_para.runs[0].font  # pragma: no cover
    f.name = font_name
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = _rgb_from_hex(color_hex)


def _add_title_bar(slide_obj, *, x, y, w, h, color_hex: str):
    bar = slide_obj.shapes.add_shape(1, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb_from_hex(color_hex)
    bar.line.fill.background()


def _add_footer(slide_obj, *, theme: DeckTheme, slide_number: int, page_w, page_h, deck_label: str = "Data Quantum"):
    from pptx.util import Inches

    margin_l = Inches(theme.spacing.get("margin_left_in", 0.7))
    margin_r = Inches(theme.spacing.get("margin_right_in", 0.7))
    footer_h = Inches(theme.spacing.get("footer_height_in", 0.25))
    y = page_h - footer_h - Inches(0.05)

    left_box = slide_obj.shapes.add_textbox(margin_l, y, page_w / 2, footer_h)
    tf_left = left_box.text_frame
    tf_left.clear()
    p_left = tf_left.paragraphs[0]
    p_left.text = deck_label
    _apply_font(
        p_left,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("footer_pt", 10)),
        color_hex=theme.palette.get("text_secondary", "#6B7280"),
    )

    right_box = slide_obj.shapes.add_textbox(page_w / 2, y, page_w / 2 - margin_r, footer_h)
    tf_right = right_box.text_frame
    tf_right.clear()
    p_right = tf_right.paragraphs[0]
    p_right.text = str(slide_number)
    p_right.alignment = 2  # right
    _apply_font(
        p_right,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("footer_pt", 10)),
        color_hex=theme.palette.get("text_secondary", "#6B7280"),
    )


def _render_cover(prs, *, theme: DeckTheme, title: str, subtitle: str, kicker: str):
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    _add_background(slide, width=prs.slide_width, height=prs.slide_height, color_hex=theme.palette.get("light_bg", "#F8F9FA"))
    _add_title_bar(
        slide,
        x=Inches(0),
        y=Inches(0),
        w=prs.slide_width,
        h=Inches(0.28),
        color_hex=theme.palette.get("navy", "#012061"),
    )

    kicker_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.9), Inches(11.2), Inches(0.8))
    ktf = kicker_box.text_frame
    ktf.clear()
    kp = ktf.paragraphs[0]
    kp.text = kicker or "HABILIDADES TRANSVERSALES"
    _apply_font(
        kp,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("cover_kicker_pt", 20)),
        color_hex=theme.palette.get("green", "#00CC66"),
        bold=True,
    )

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.2), Inches(2.1))
    tf_title = title_box.text_frame
    tf_title.clear()
    p = tf_title.paragraphs[0]
    p.text = title or "Curso generado"
    _apply_font(
        p,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("cover_title_pt", 44)),
        color_hex=theme.palette.get("navy", "#012061"),
        bold=True,
    )

    subtitle_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(10.8), Inches(1.7))
    tf_sub = subtitle_box.text_frame
    tf_sub.clear()
    ps = tf_sub.paragraphs[0]
    ps.text = subtitle or "Generado automaticamente para revision de administracion"
    _apply_font(
        ps,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("subtitle_pt", 20)),
        color_hex=theme.palette.get("text_secondary", "#6B7280"),
    )
    return slide


def _render_agenda(prs, *, theme: DeckTheme, items: List[str]):
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    _add_background(slide, width=prs.slide_width, height=prs.slide_height, color_hex=theme.palette.get("white", "#FFFFFF"))
    _add_title_bar(slide, x=Inches(0.7), y=Inches(0.65), w=Inches(1.3), h=Inches(0.18), color_hex=theme.palette.get("green", "#00CC66"))

    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(5.0), Inches(0.8))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Indice de Modulos"
    _apply_font(
        p,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("slide_title_pt", 30)),
        color_hex=theme.palette.get("navy", "#012061"),
        bold=True,
    )

    list_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(10.6), Inches(4.7))
    tf_list = list_box.text_frame
    tf_list.clear()
    for i, item in enumerate(items[:8]):
        para = tf_list.paragraphs[0] if i == 0 else tf_list.add_paragraph()
        para.text = f"{i + 1}  {item}"
        _apply_font(
            para,
            font_name=theme.typography.get("font_family", "Calibri"),
            size_pt=int(theme.typography.get("body_pt", 18)),
            color_hex=theme.palette.get("text_primary", "#0F172A"),
        )
    return slide


def _render_section_divider(prs, *, theme: DeckTheme, title: str, subtitle: str, module_number: int):
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    _add_background(slide, width=prs.slide_width, height=prs.slide_height, color_hex=theme.palette.get("navy", "#012061"))
    _add_title_bar(slide, x=Inches(0), y=Inches(0), w=Inches(12.8), h=Inches(0.24), color_hex=theme.palette.get("green", "#00CC66"))

    mbox = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(4.0), Inches(0.9))
    mtf = mbox.text_frame
    mtf.clear()
    mp = mtf.paragraphs[0]
    mp.text = f"MODULO {module_number}"
    _apply_font(
        mp,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=20,
        color_hex=theme.palette.get("green", "#00CC66"),
        bold=True,
    )

    tbox = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.0), Inches(1.4))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title or "Seccion"
    _apply_font(
        p,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("slide_title_pt", 30)),
        color_hex=theme.palette.get("white", "#FFFFFF"),
        bold=True,
    )

    sbox = slide.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(10.8), Inches(0.9))
    stf = sbox.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = subtitle or "Bloque de contenido"
    _apply_font(
        sp,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("subtitle_pt", 20)),
        color_hex=theme.palette.get("green", "#00CC66"),
    )
    return slide


def _render_content_bullets(
    prs,
    *,
    theme: DeckTheme,
    title: str,
    bullets: List[str],
    script: str,
    module_label: str,
    progress_label: str,
):
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    _add_background(slide, width=prs.slide_width, height=prs.slide_height, color_hex=theme.palette.get("white", "#FFFFFF"))
    _add_title_bar(slide, x=Inches(0.7), y=Inches(0.6), w=Inches(2.4), h=Inches(0.18), color_hex=theme.palette.get("navy", "#012061"))

    module_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.82), Inches(8.0), Inches(0.5))
    mtf = module_box.text_frame
    mtf.clear()
    mp = mtf.paragraphs[0]
    mp.text = module_label
    _apply_font(
        mp,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=14,
        color_hex=theme.palette.get("text_secondary", "#6B7280"),
    )

    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(11.2), Inches(0.9))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    _apply_font(
        p,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("slide_title_pt", 30)),
        color_hex=theme.palette.get("navy", "#012061"),
        bold=True,
    )

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(10.7), Inches(4.1))
    btf = body.text_frame
    btf.clear()
    for i, b in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = b
        para.level = 0
        _apply_font(
            para,
            font_name=theme.typography.get("font_family", "Calibri"),
            size_pt=int(theme.typography.get("body_pt", 18)),
            color_hex=theme.palette.get("text_primary", "#0F172A"),
        )

    if script:
        note = slide.shapes.add_textbox(Inches(0.9), Inches(6.15), Inches(10.7), Inches(0.45))
        ntf = note.text_frame
        ntf.clear()
        np = ntf.paragraphs[0]
        np.text = f"Nota ponente: {script[:190]}"
        _apply_font(
            np,
            font_name=theme.typography.get("font_family", "Calibri"),
            size_pt=11,
            color_hex=theme.palette.get("text_secondary", "#6B7280"),
        )
    pbox = slide.shapes.add_textbox(Inches(10.5), Inches(6.15), Inches(1.2), Inches(0.45))
    ptf = pbox.text_frame
    ptf.clear()
    pp = ptf.paragraphs[0]
    pp.text = progress_label
    pp.alignment = 2
    _apply_font(
        pp,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=11,
        color_hex=theme.palette.get("text_secondary", "#6B7280"),
    )
    return slide


def _render_content_two_columns(prs, *, theme: DeckTheme, title: str, bullets: List[str], script: str):
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    _add_background(slide, width=prs.slide_width, height=prs.slide_height, color_hex=theme.palette.get("light_bg", "#F8F9FA"))
    _add_title_bar(slide, x=Inches(0.7), y=Inches(0.6), w=Inches(2.6), h=Inches(0.18), color_hex=theme.palette.get("orange", "#F97316"))

    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(11.2), Inches(0.9))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    _apply_font(
        p,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("slide_title_pt", 30)),
        color_hex=theme.palette.get("navy", "#012061"),
        bold=True,
    )

    half = max(1, (len(bullets) + 1) // 2)
    left_items = bullets[:half]
    right_items = bullets[half:]

    lbox = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(5.1), Inches(4.2))
    ltf = lbox.text_frame
    ltf.clear()
    for i, b in enumerate(left_items):
        para = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        para.text = b
        _apply_font(
            para,
            font_name=theme.typography.get("font_family", "Calibri"),
            size_pt=int(theme.typography.get("body_pt", 18)),
            color_hex=theme.palette.get("text_primary", "#0F172A"),
        )

    rbox = slide.shapes.add_textbox(Inches(6.4), Inches(1.9), Inches(5.1), Inches(4.2))
    rtf = rbox.text_frame
    rtf.clear()
    for i, b in enumerate(right_items):
        para = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        para.text = b
        _apply_font(
            para,
            font_name=theme.typography.get("font_family", "Calibri"),
            size_pt=int(theme.typography.get("body_pt", 18)),
            color_hex=theme.palette.get("text_primary", "#0F172A"),
        )

    if script:
        note = slide.shapes.add_textbox(Inches(0.9), Inches(6.15), Inches(10.7), Inches(0.45))
        ntf = note.text_frame
        ntf.clear()
        np = ntf.paragraphs[0]
        np.text = f"Nota ponente: {script[:160]}"
        _apply_font(
            np,
            font_name=theme.typography.get("font_family", "Calibri"),
            size_pt=11,
            color_hex=theme.palette.get("text_secondary", "#6B7280"),
        )
    return slide


def _render_closing(prs, *, theme: DeckTheme, title: str, next_steps: List[str]):
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    _add_background(slide, width=prs.slide_width, height=prs.slide_height, color_hex=theme.palette.get("white", "#FFFFFF"))
    _add_title_bar(slide, x=Inches(0), y=Inches(0), w=Inches(12.8), h=Inches(0.24), color_hex=theme.palette.get("navy", "#012061"))

    tbox = slide.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(10.8), Inches(0.9))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title or "Proximos pasos"
    _apply_font(
        p,
        font_name=theme.typography.get("font_family", "Calibri"),
        size_pt=int(theme.typography.get("slide_title_pt", 30)),
        color_hex=theme.palette.get("navy", "#012061"),
        bold=True,
    )

    body = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(10.5), Inches(3.8))
    btf = body.text_frame
    btf.clear()
    for i, b in enumerate(next_steps[:6]):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = f"- {b}"
        _apply_font(
            para,
            font_name=theme.typography.get("font_family", "Calibri"),
            size_pt=int(theme.typography.get("body_pt", 18)),
            color_hex=theme.palette.get("text_primary", "#0F172A"),
        )
    return slide


def _build_semantic_blocks(
    slides: List[Dict[str, Any]],
    *,
    max_bullets_per_slide: int,
    max_chars_per_bullet: int,
) -> Tuple[List[Dict[str, Any]], int, int]:
    """
    Convierte slides semanticas en bloques renderizables y aplica split/truncate.
    """
    blocks: List[Dict[str, Any]] = []
    trunc_count = 0
    split_count = 0

    for i, slide in enumerate(slides, start=1):
        title = _safe_title(slide, i)
        bullets = slide.get("bullets") or []
        script = (slide.get("script") or "").strip()
        chunks, c_trunc = _split_bullets(
            bullets,
            max_bullets=max_bullets_per_slide,
            max_chars=max_chars_per_bullet,
        )
        trunc_count += c_trunc
        if len(chunks) > 1:
            split_count += (len(chunks) - 1)

        for chunk_idx, chunk in enumerate(chunks):
            suffix = f" (Parte {chunk_idx + 1})" if len(chunks) > 1 else ""
            blocks.append(
                {
                    "title": f"{title}{suffix}",
                    "bullets": chunk,
                    "script": script,
                    "origin_slide": i,
                }
            )
    return blocks, trunc_count, split_count


def build_pptx_deck(
    *,
    course_id: str,
    slides: List[Dict[str, Any]],
    images_by_slide_number: List[Dict[str, Any]],  # noqa: ARG001 - reservado para futuras versiones
    deck_title: Optional[str] = None,
    course_description: Optional[str] = None,
    learning_objectives: Optional[List[str]] = None,
) -> Tuple[str, str | None]:
    """
    Renderiza un deck con preset corporativo fijo y contenido del LLM.

    Contract de salida estable:
      (deck_file_url, err)
    """
    try:
        from pptx import Presentation
    except Exception as e:  # pragma: no cover
        return "", f"Falta python-pptx/Pillow en runtime: {str(e)}"

    generation_dir = os.path.join(_static_generated_root(), course_id)
    os.makedirs(generation_dir, exist_ok=True)

    try:
        theme = _load_theme()
    except Exception as e:
        return "", f"No se pudo cargar el theme preset: {str(e)}"

    prs = Presentation()
    max_bullets = int(theme.limits.get("max_bullets_per_slide", 6))
    max_chars = int(theme.limits.get("max_chars_per_bullet", 110))

    semantic_blocks, trunc_count, split_count = _build_semantic_blocks(
        slides,
        max_bullets_per_slide=max_bullets,
        max_chars_per_bullet=max_chars,
    )

    layout_order = theme.layout_order or [
        "cover",
        "agenda",
        "section-divider",
        "content-bullets",
        "content-two-columns",
        "closing",
    ]
    content_layouts = [l for l in layout_order if l in ("content-bullets", "content-two-columns")]
    if not content_layouts:
        content_layouts = ["content-bullets", "content-two-columns"]

    module_size = int(theme.structure.get("module_size", 3))
    modules: List[List[Dict[str, Any]]] = []
    for i in range(0, len(semantic_blocks), module_size):
        modules.append(semantic_blocks[i : i + module_size])

    # 1) Cover
    cover = _render_cover(
        prs,
        theme=theme,
        title=(deck_title or "Curso generado"),
        subtitle=f"{course_description or 'Ruta formativa para revision de administracion'}\nNivel: {('N/A')}",
        kicker="HABILIDADES TRANSVERSALES",
    )
    _add_footer(
        cover,
        theme=theme,
        slide_number=1,
        page_w=prs.slide_width,
        page_h=prs.slide_height,
        deck_label="dataquantum",
    )
    print("[PPTX] layout seleccionado: cover", flush=True)

    # 2) Agenda
    agenda_items = [m[0].get("title", "") for m in modules[:8] if m]
    agenda = _render_agenda(prs, theme=theme, items=agenda_items)
    _add_footer(
        agenda,
        theme=theme,
        slide_number=2,
        page_w=prs.slide_width,
        page_h=prs.slide_height,
        deck_label=f"dataquantum  ·  {deck_title or 'Curso'}",
    )
    print("[PPTX] layout seleccionado: agenda", flush=True)

    slide_counter = 3

    # 3) Divisor + contenido por modulo
    for mod_idx, module_blocks in enumerate(modules, start=1):
        divider = _render_section_divider(
            prs,
            theme=theme,
            title=module_blocks[0].get("title") if module_blocks else "Contenido",
            subtitle="Bloque de aprendizaje",
            module_number=mod_idx,
        )
        _add_footer(
            divider,
            theme=theme,
            slide_number=slide_counter,
            page_w=prs.slide_width,
            page_h=prs.slide_height,
            deck_label=f"dataquantum  ·  {deck_title or 'Curso'}",
        )
        print("[PPTX] layout seleccionado: section-divider", flush=True)
        slide_counter += 1

        total_in_module = max(1, len(module_blocks))
        for idx, block in enumerate(module_blocks):
            layout_type = content_layouts[idx % len(content_layouts)]
            module_label = f"Modulo {mod_idx}: {module_blocks[0].get('title','Contenido')}"
            progress_label = f"{idx + 1} / {total_in_module}"
            if layout_type == "content-bullets":
                s = _render_content_bullets(
                    prs,
                    theme=theme,
                    title=block.get("title") or "Contenido",
                    bullets=block.get("bullets") or [],
                    script=block.get("script") or "",
                    module_label=module_label,
                    progress_label=progress_label,
                )
            else:
                s = _render_content_two_columns(
                    prs,
                    theme=theme,
                    title=block.get("title") or "Contenido",
                    bullets=block.get("bullets") or [],
                    script=block.get("script") or "",
                )
            _add_footer(
                s,
                theme=theme,
                slide_number=slide_counter,
                page_w=prs.slide_width,
                page_h=prs.slide_height,
                deck_label=f"dataquantum  ·  {deck_title or 'Curso'}",
            )
            print(f"[PPTX] layout seleccionado: {layout_type} | title={block.get('title','')[:60]}", flush=True)
            slide_counter += 1

    # 5) Closing
    next_steps = (learning_objectives or [])[:6]
    if not next_steps:
        next_steps = [
            "Revisar los conceptos principales",
            "Aplicar el contenido en un caso real",
            "Compartir resultados y mejoras con el equipo",
        ]
    closing = _render_closing(
        prs,
        theme=theme,
        title="Proximos pasos",
        next_steps=next_steps,
    )
    _add_footer(closing, theme=theme, slide_number=slide_counter, page_w=prs.slide_width, page_h=prs.slide_height)
    print("[PPTX] layout seleccionado: closing", flush=True)

    if trunc_count > 0:
        print(f"[PPTX] truncados aplicados: {trunc_count}", flush=True)
    if split_count > 0:
        print(f"[PPTX] splits aplicados: {split_count}", flush=True)

    deck_file_name = "deck.pptx"
    deck_file_path = os.path.join(generation_dir, deck_file_name)
    prs.save(deck_file_path)
    print(f"[PPTX] ruta final del deck: {deck_file_path}", flush=True)

    deck_file_url = f"/static/generated-courses/{course_id}/{deck_file_name}"
    return deck_file_url, None

